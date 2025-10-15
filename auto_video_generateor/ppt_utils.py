# pip install python-pptx pdf2image PyMuPDF
import os
import shutil
import subprocess
from pptx import Presentation
import fitz  # PyMuPDF

from common_utils import get_savepath, chat


def ppt_to_pdf(inpath, code_name):
    """需要下载安装LibreOffice软件，培训环境变量，自行搜索解决安装问题。"""
    pdf_dir = get_savepath(code_name, 'pdf', mkdir_ok=True)
    if inpath.endswith('.pdf'):
        pdf_path = os.path.join(pdf_dir, 'document.pdf')
        shutil.copyfile(inpath, pdf_path)
        ppt_path = ''
        return ppt_path, pdf_path

    ppt_path = os.path.join(pdf_dir, 'document.pptx')
    shutil.copyfile(inpath, ppt_path)

    # Step 1: 使用libreoffice将PPT转换为PDF
    # libreoffice_path = r"C:\Program Files\LibreOffice\program\soffice.exe"
    libreoffice_path = "soffice"

    pdf_path = os.path.join(pdf_dir, os.path.splitext(os.path.basename(ppt_path))[0] + ".pdf")
    command = [
        libreoffice_path,
        "--headless",
        "--convert-to", "pdf",
        ppt_path,
        "--outdir", pdf_dir
    ]

    # 执行转换命令
    subprocess.run(command, check=True)
    print(f"PDF saved at: {pdf_path}")
    return ppt_path, pdf_path


def pdf_to_images(pdf_path, ij_dict, code_name):
    image_dir = get_savepath(code_name, 'image', mkdir_ok=True)

    # Step 2: 将PDF的每一页转换为图片
    image_format = "png"
    dpi = 150
    # 检查 PDF 文件是否存在
    if not os.path.isfile(pdf_path):
        raise FileNotFoundError(f"The file {pdf_path} does not exist.")

    # 打开 PDF 文件
    pdf_document = fitz.open(pdf_path)
    images = []
    # 遍历每一页
    j_cnt = 0
    for page_num in range(pdf_document.page_count):
        # 获取页面
        page = pdf_document[page_num]

        # 设置缩放比例
        zoom = dpi / 72  # 默认 PDF 分辨率是 72 DPI
        matrix = fitz.Matrix(zoom, zoom)

        # 渲染页面为图片
        pix = page.get_pixmap(matrix=matrix)
        for j in range(ij_dict[page_num]):
            # 保存图片
            image_path = f'{image_dir}/image_{j_cnt + 100}.{image_format}'
            pix.save(image_path)
            print(f"Saved: {image_path}")
            images.append(image_path)
            j_cnt += 1
    return images


def pdf_to_texts(pdf_path, prompt_template, code_name):
    # text_dir = get_savepath(code_name, 'text', mkdir_ok=True)

    # Step 2: 将PDF的每一页转换为文字
    # 检查 PDF 文件是否存在
    if not os.path.isfile(pdf_path):
        raise FileNotFoundError(f"The file {pdf_path} does not exist.")

    # 打开 PDF 文件
    pdf_document = fitz.open(pdf_path)
    texts = []
    # 遍历每一页
    for page_num in range(pdf_document.page_count):
        # 获取页面
        page = pdf_document[page_num]

        text = page.get_text()  # 提取当前页的文本
        prompt_text = prompt_template.format(text)
        note_text = chat(prompt_text)
        texts.append(note_text)
        # note_path = f'{text_dir}/text_{page_num + 100}.txt'
        # with open(note_path, 'w', encoding='utf-8') as file:
        #     file.write(note_text)
        #
        # print(f"Saved notes: {note_path}")
    return texts


def ppt_to_texts(ppt_path, code_name):
    # text_dir = get_savepath(code_name, 'text', mkdir_ok=True)

    # Step 3: 使用python-pptx提取每页备注并保存为TXT文件
    ppt = Presentation(ppt_path)
    texts = []
    for i, slide in enumerate(ppt.slides):
        note_text = slide.notes_slide.notes_text_frame.text if slide.has_notes_slide else ""

        # 遍历幻灯片中的所有形状
        # texts_tmp = []
        # for shape in slide.shapes:
        #     if hasattr(shape, 'text'):
        #         # 提取形状中的文字
        #         texts_tmp.append(shape.text)
        # texts_tmp = '。'.join(texts_tmp)
        # note_text = note_text or texts_tmp or "嗯。"

        texts.append(note_text)
        # note_path = f'{text_dir}/text_{i + 100}.txt'
        # with open(note_path, 'w', encoding='utf-8') as file:
        #     file.write(note_text)
        #
        # print(f"Saved notes: {note_path}")
    return texts
